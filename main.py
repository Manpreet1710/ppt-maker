from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Dict, Optional
import google.generativeai as genai
from pptx import Presentation
import io
from datetime import datetime
import json
from pathlib import Path
from typing import List, Dict, Optional, Any
import os
from pptx.util import Pt


app = FastAPI(title="AI PPT Generator API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # यह सभी origins से अनुरोधों की अनुमति देता है
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"]
)

GEMINI_API_KEY = "AIzaSyBy0NtAiwd7wkSHNKreNpWBEoJ68yi8uFI"
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-2.0-flash")

TEMPLATE_DIR = Path("templates")
TEMPLATE_DIR.mkdir(exist_ok=True)


class ContentGenRequest(BaseModel):
    topic: str
    slide_count: int = 6
    language: str = "English"
    tone: str = "Professional"


class GeneratePPTRequest(BaseModel):
    template_id: str
    content: Any


TEMPLATE_INFO = {
    "template-1.pptx": {
        "id": "template_1",
        "name": "Blue Minimalist Report",
        "category": "Business",
        "slides_count": 5,
        "thumbnail_url": "/templates/previews/template-1.png"
    },
    "template-2.pptx": {
        "id": "template_2",
        "name": "Dark Professional Pitch",
        "category": "Pitch Deck",
        "slides_count": 5,
        "thumbnail_url": "/templates/previews/template-2.png"
    },
    "template-3.pptx": {
        "id": "template_3",
        "name": "Light Corporate",
        "category": "Corporate",
        "slides_count": 5,
        "thumbnail_url": "/templates/previews/template-3.png"
    }
}

# -----------------------------------------------------------------------------
# API Endpoints
# -----------------------------------------------------------------------------


@app.get("/")
def root():
    return {
        "message": "AI PPT Generator API",
        "version": "4.0",
        "status": "running",
        "endpoints": {
            "generate_content": "POST /generate-content",
            "generate_ppt": "POST /generate-ppt",
            "templates": "GET /templates",
            "health": "GET /health"
        }
    }


@app.get("/health")
def health_check():
    """Check system health"""
    templates_status = {}
    for filename in TEMPLATE_INFO.keys():
        path = TEMPLATE_DIR / filename
        templates_status[filename] = "found" if path.exists() else "missing"

    return {
        "status": "healthy",
        "gemini_configured": bool(GEMINI_API_KEY),
        "templates": templates_status
    }


@app.get("/templates")
def list_templates():
    """Get available templates"""
    templates = []
    for filename, info in TEMPLATE_INFO.items():
        templates.append({
            "id": info["id"],
            "name": info["name"],
            "category": info["category"],
            "slides_count": info["slides_count"],
            "thumbnail_url": info["thumbnail_url"]
        })

    return {"templates": templates}



@app.post("/generate-content")
async def generate_content(request: ContentGenRequest):
    """Generate hierarchical content using Gemini AI"""
    if not request.topic.strip():
        raise HTTPException(status_code=400, detail="Topic cannot be empty")

    prompt = f"""
You are an expert presentation content generator.

TASK:
Create a detailed, hierarchical JSON structure for: "{request.topic}"

LANGUAGE: {request.language}
TONE: {request.tone}
SLIDES: {request.slide_count}

STRICT JSON OUTPUT RULES:
1. Return ONLY valid JSON - no markdown, no code blocks, no explanations
2. Do NOT use ``` or ```json
3. Start with {{ and end with }}
4. Follow the EXACT structure below

REQUIRED JSON STRUCTURE:
{{
  "data": {{
    "level": 1,
    "name": "A Comprehensive Exploration",
    "type": null,
    "children": [
      {{
        "level": 2,
        "name": "Section Title",
        "url": null,
        "type": null,
        "children": [...]
      }}
    ]
  }}
}}

CONTENT GENERATION RULES:
1. DO NOT create a section titled "Introduction" - start directly with meaningful content sections
2. Create EXACTLY {request.slide_count} main sections (level 2 nodes) based on the topic
3. Analyze the topic and generate relevant, specific section titles (NOT generic like "Introduction", "Conclusion")
4. Each section must have 2-3 subsections (level 3)
5. Each subsection must have 2-3 key points (level 4)
6. Each key point must have 1 detailed explanation (level 0)
7. Level 0 content should be 1-2 complete sentences with rich context
8. Use professional, informative language
9. Maintain logical flow throughout the content
10. Keep all "url" and "type" as null
11. All "children" arrays must exist (use [] if empty)

SECTION NAMING RULES:
- ❌ AVOID: "Introduction", "Conclusion", "Overview", "Summary"
- ✅ USE: Specific, topic-relevant titles based on the content
- Example for "Climate Change": "Key Indicators", "Causes", "Impacts", "Mitigation", "Future Outlook", "Global Response"
- Example for "Artificial Intelligence": "Machine Learning Basics", "Neural Networks", "AI Applications", "Ethical Considerations", "Future Trends", "Industry Impact"

CONTENT QUALITY:
- Level 2: Short, concise section titles (2-4 words max) - specific to the topic
- Level 3: Specific subsection topics
- Level 4: Focused key points
- Level 0: Rich, detailed explanations with context and insights (1-2 sentences)

OUTPUT:
Return ONLY the JSON starting with {{ and ending with }}. No other text.
"""
    try:
        # Streaming ko collect karo
        full_text = ""
        stream = model.generate_content(prompt, stream=True)
        
        for chunk in stream:
            if chunk.text:
                full_text += chunk.text
        
        # Clean the response
        full_text = full_text.strip()
        
        # Remove markdown code blocks if present
        if full_text.startswith('```'):
            full_text = full_text.split('```')[1]
            if full_text.startswith('json'):
                full_text = full_text[4:]
            full_text = full_text.strip()
        
        # Parse JSON
        data = json.loads(full_text)
        
        # Return complete response
        return {
            "status": "success",
            "data": data
        }

    except json.JSONDecodeError as e:
        return {
            "status": "error",
            "error": f"Invalid JSON: {str(e)}",
            "raw_response": full_text[:500]  # First 500 chars for debugging
        }
    except Exception as e:
        return {
            "status": "error",
            "error": str(e)
        }


def get_template_path(template_id: str) -> str:
    """Get template file path from template_id"""
    for filename, info in TEMPLATE_INFO.items():
        if info["id"] == template_id:
            path = TEMPLATE_DIR / filename
            if not path.exists():
                raise FileNotFoundError(f"Template file not found: {filename}")
            return str(path)
    raise ValueError(f"Invalid template_id: {template_id}")

def extract_template_placeholders(template_path: str):
    prs = Presentation(template_path)
    template_structure = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": slide_index,
            "placeholders": []
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_info["placeholders"].append({
                    "id": id(shape),
                    "type": "BODY" if shape.is_placeholder else "TEXT",
                    "text": shape.text
                })
        template_structure.append(slide_info)

    return template_structure


def set_text_preserve_format(shape, new_text):
    """Replace text in a shape while preserving formatting."""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    
    # Clear existing text
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = ""
    
    # Add new text with preserved formatting
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        first_run = text_frame.paragraphs[0].runs[0]
        new_run = text_frame.paragraphs[0].add_run()
        new_run.text = new_text
        
        # Copy formatting
        if first_run.font.name:
            new_run.font.name = first_run.font.name
        if first_run.font.size:
            new_run.font.size = first_run.font.size
        new_run.font.bold = first_run.font.bold
        new_run.font.italic = first_run.font.italic
        if first_run.font.color.rgb:
            new_run.font.color.rgb = first_run.font.color.rgb
    else:
        text_frame.text = new_text

def fill_template_with_ai_content(prs: Presentation, title: str, children: list) -> Presentation:
    """
    Fill PPT template with AI-generated content.
    
    Template Structure (21 slides):
    - Slide 1: Title
    - Slide 2: Agenda
    - Slides 3-21: 6 sections, each with:
        * Section divider (title + number)
        * 2 content slides per section
    """
    print("\n===== FILLING PPT CONTENT =====")
    
    slides = prs.slides
    
    # ========================================
    # SLIDE 1: TITLE SLIDE
    # ========================================
    print(f"\n[Slide 1] Title = {title}")
    title_shapes = [s for s in slides[0].shapes if s.has_text_frame]
    
    if title_shapes:
        set_text_preserve_format(title_shapes[0], title)
    
    # ========================================
    # SLIDE 2: AGENDA SLIDE
    # ========================================
    agenda_slide = slides[1]
    text_shapes = [s for s in agenda_slide.shapes if s.has_text_frame]
    
    # Separate heading from items
    agenda_title_shape = None
    agenda_item_shapes = []
    
    for shape in text_shapes:
        text = shape.text.strip().upper()
        if text == "AGENDA":
            agenda_title_shape = shape
        else:
            agenda_item_shapes.append(shape)
    
    # Sort by vertical position
    agenda_item_shapes.sort(key=lambda x: x.top)
    
    # Fill agenda items with Level-2 names from AI content
    agenda_items = [child["name"] for child in children]
    
    print(f"\n[Slide 2] Agenda Items:")
    for i, item in enumerate(agenda_items):
        print(f"   {i+1}. {item}")
        if i < len(agenda_item_shapes):
            set_text_preserve_format(agenda_item_shapes[i], item)
    
    # ========================================
    # SLIDES 3-21: SECTION CONTENT
    # ========================================
    # Pattern: Every section has 3 slides (divider + 2 content)
    # Slide indices: 2, 5, 8, 11, 14, 17 (dividers for 6 sections)
    
    section_divider_slides = [2, 6, 9, 12, 15, 18]  # 0-indexed: slides 3, 7, 10, 13, 16, 19
    
    print("\n=== PROCESSING SECTIONS ===")
    
    for section_idx, child in enumerate(children):
        if section_idx >= 6:  # Template only has 6 sections
            break
        
        # Get section divider slide index
        divider_idx = section_divider_slides[section_idx]
        
        if divider_idx >= len(slides):
            break
        
        print(f"\n--- SECTION {section_idx + 1}: {child['name']} ---")
        
        # ========================================
        # SECTION DIVIDER SLIDE
        # ========================================
        divider_slide = slides[divider_idx]
        divider_shapes = [s for s in divider_slide.shapes if s.has_text_frame]
        
        # Separate title shape from number shape
        title_shape = None
        number_shape = None
        
        for shape in divider_shapes:
            text = shape.text.strip()
            # Check if this is the number (01, 02, etc.)
            if text in ["01", "02", "03", "04", "05", "06"]:
                number_shape = shape
            else:
                title_shape = shape
        
        # Fill section title (DO NOT change the number)
        if title_shape:
            print(f"[Slide {divider_idx + 1}] Section Title: {child['name']}")
            set_text_preserve_format(title_shape, child['name'])
        
        # Keep the number unchanged
        if number_shape:
            print(f"[Slide {divider_idx + 1}] Section Number: {number_shape.text} (preserved)")
            # Don't modify number_shape at all
        
        # ========================================
        # CONTENT SLIDES (2 slides after divider)
        # ========================================
        level3_items = child.get("children", [])
        content_slide_indices = [divider_idx + 1, divider_idx + 2]
        
        content_pointer = 0  # Track position in level3_items
        
        for content_slide_idx in content_slide_indices:
            if content_slide_idx >= len(slides):
                break
            
            content_slide = slides[content_slide_idx]
            content_shapes = [s for s in content_slide.shapes if s.has_text_frame]
            
            # Sort by position: top-to-bottom within left-to-right columns
            # Divide slide into left/right halves based on shape positions
            # Default PowerPoint slide width is ~9144000 EMUs (10 inches)
            
            if content_shapes:
                # Find the midpoint based on actual shape positions
                all_lefts = [s.left for s in content_shapes]
                min_left = min(all_lefts)
                max_left = max(all_lefts)
                mid_point = (min_left + max_left) / 2
            else:
                mid_point = 4572000  # Default half of standard slide width
            
            left_shapes = [s for s in content_shapes if s.left < mid_point]
            right_shapes = [s for s in content_shapes if s.left >= mid_point]
            
            left_shapes.sort(key=lambda x: x.top)
            right_shapes.sort(key=lambda x: x.top)
            
            # Combine: all left shapes first, then right shapes
            sorted_shapes = left_shapes + right_shapes
            
            print(f"\n[Slide {content_slide_idx + 1}] Layout:")
            print(f"   Left column: {len(left_shapes)} shapes")
            print(f"   Right column: {len(right_shapes)} shapes")
            print(f"   Fill order: Left-top to bottom, then Right-top to bottom")
            
            shape_idx = 0
            
            # Fill this slide with level-3 and level-4 content
            while content_pointer < len(level3_items) and shape_idx < len(sorted_shapes):
                level3 = level3_items[content_pointer]
                level4_items = level3.get("children", [])
                
                # LEVEL-3 Title (e.g., "Defining Love: A Multifaceted Emotion")
                print(f"\n   📌 L3: {level3['name']}")
                if shape_idx < len(sorted_shapes):
                    set_text_preserve_format(sorted_shapes[shape_idx], level3['name'])
                    print(f"      → Filled shape {shape_idx + 1}")
                    shape_idx += 1
                
                # LEVEL-4 Items with their level-0 descriptions
                for level4 in level4_items:
                    if shape_idx >= len(sorted_shapes):
                        break
                    
                    # LEVEL-4 Heading (e.g., "The Spectrum of Love")
                    print(f"   📍 L4: {level4['name']}")
                    set_text_preserve_format(sorted_shapes[shape_idx], level4['name'])
                    print(f"      → Filled shape {shape_idx + 1}")
                    shape_idx += 1
                    
                    # LEVEL-0 Description
                    level0_items = level4.get("children", [])
                    if level0_items and shape_idx < len(sorted_shapes):
                        description = level0_items[0]["name"]
                        print(f"   📄 L0: {description[:60]}...")
                        set_text_preserve_format(sorted_shapes[shape_idx], description)
                        print(f"      → Filled shape {shape_idx + 1}")
                        shape_idx += 1
                
                content_pointer += 1
                
                # If we've filled all shapes on this slide, move to next slide
                if shape_idx >= len(sorted_shapes):
                    print(f"   ✅ Slide full, moving to next slide")
                    break
    
    print("\n===== PPT BUILD COMPLETE =====\n")
    return prs

# Main execution
def create_presentation(template_path: str, ai_content: dict):
    """Main function to create presentation."""
    
    # Extract data
    title = ai_content["data"].get("name", "Presentation Title")
    children = ai_content["data"].get("children", [])
    
    # Load template
    prs = Presentation(template_path)
    
    # Fill content
    prs = fill_template_with_ai_content(prs, title, children)
    
    # Save
    output_path = Path("output_ppt.pptx")
    prs.save(output_path)
    
    return {
        "status": "success",
        "message": "PPT generated successfully",
        "output_file": str(output_path)
    }



import subprocess
from pathlib import Path
from fastapi.responses import FileResponse

def convert_ppt_to_pdf(pptx_path: Path) -> Path:
    pdf_path = pptx_path.with_suffix('.pdf')
    
    try:
        result = subprocess.run([
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '--headless',
            '--invisible',
            '--nocrashreport',
            '--nodefault',
            '--nofirststartwizard',
            '--nolockcheck',
            '--nologo',
            '--norestore',
            '--convert-to', 'pdf',
            '--outdir', str(pptx_path.parent),
            str(pptx_path)
        ], check=True, timeout=60, capture_output=True, text=True)
        
        if not pdf_path.exists():
            raise Exception(f"PDF not created. Output: {result.stdout}")
        
        return pdf_path
    except subprocess.TimeoutExpired:
        raise Exception("PDF conversion timeout (60s)")
    except FileNotFoundError:
        raise Exception("LibreOffice not found. Install: brew install --cask libreoffice")
    except Exception as e:
        raise Exception(f"PDF conversion failed: {e}")


import io

@app.post("/create-presentation")
async def create_presentation(request: GeneratePPTRequest):
    ai_content = request.content
    title = ai_content["data"].get("name", "")
    children = ai_content["data"].get("children", [])

    template_path = get_template_path(request.template_id)
    prs = Presentation(template_path)
    prs = fill_template_with_ai_content(prs, title, children)

    # In-memory buffer mein save karo (RAM)
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)  # Start se read karne ke liye

    # PDF bhi in-memory (optional - agar chahiye to)
    pdf_buffer = None
    try:
        # Temporary file sirf PDF conversion ke liye
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
            temp_pptx.write(pptx_buffer.getvalue())
            temp_pptx_path = temp_pptx.name
        
        # Convert to PDF
        pdf_path = convert_ppt_to_pdf(Path(temp_pptx_path))
        
        # PDF ko bhi memory mein load karo
        with open(pdf_path, 'rb') as f:
            pdf_buffer = io.BytesIO(f.read())
            pdf_buffer.seek(0)
        
        # Cleanup temp files
        Path(temp_pptx_path).unlink()
        pdf_path.unlink()
        
    except Exception as e:
        print(f"⚠️ PDF conversion failed: {e}")

    # Base64 encode karke return karo (ya session mein store karo)
    import base64
    pptx_base64 = base64.b64encode(pptx_buffer.getvalue()).decode()
    pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode() if pdf_buffer else None

    return {
        "status": "success",
        "message": "PPT generated successfully",
        "pptx_data": pptx_base64,  # Frontend download karega
        "pdf_data": pdf_base64
    }


@app.get("/download-pptx/{file_data}")
async def download_pptx(file_data: str):
    """Direct download from base64 data"""
    import base64
    
    # Base64 decode karo
    file_bytes = base64.b64decode(file_data)
    buffer = io.BytesIO(file_bytes)
    
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"}
    )


@app.get("/download-pdf/{file_data}")
async def download_pdf(file_data: str):
    """Direct download from base64 data"""
    import base64
    
    file_bytes = base64.b64decode(file_data)
    buffer = io.BytesIO(file_bytes)
    
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": "inline; filename=presentation.pdf"}
    )

@app.get("/debug-template")
def debug_template(template_id: str):
    template_path = get_template_path(template_id)
    prs = Presentation(template_path)

    info = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {"slide_number": s_idx, "shapes": []}
        
        # टेक्स्ट वाली shapes के लिए एक अलग counter शुरू करें
        text_shape_counter = 0 

        # सभी shapes पर iterate करें
        for shape in slide.shapes:
            if shape.has_text_frame:
                # टेक्स्ट वाली shape मिलने पर counter बढ़ाएँ
                text_shape_counter += 1 
                
                paragraphs = [p.text for p in shape.text_frame.paragraphs]
                text = "\n".join(paragraphs)

                slide_data["shapes"].append({
                    # अब text_shape_counter का उपयोग करें 
                    "shape_index": text_shape_counter, 
                    "placeholder": shape.is_placeholder,
                    "text": text
                })

        info.append(slide_data)

    return info




if __name__ == "__main__":
    import uvicorn

    print("=" * 60)
    print("🚀 AI PPT Generator API Server")
    print("=" * 60)
    print(f"📁 Templates directory: {TEMPLATE_DIR.absolute()}")
    print(f"📋 Available templates: {len(TEMPLATE_INFO)}")

    # Check templates
    missing = [f for f in TEMPLATE_INFO.keys() if not (
        TEMPLATE_DIR / f).exists()]
    if missing:
        print(f"⚠️  WARNING: Missing templates: {', '.join(missing)}")
        print(f"   Add these files to: {TEMPLATE_DIR.absolute()}")
    else:
        print("✅ All templates found!")

    print("=" * 60)
    print("\n📡 API Endpoints:")
    print("   POST /generate-content - Generate AI content")
    print("   POST /generate-ppt     - Create PowerPoint")
    print("   GET  /templates        - List templates")
    print("   GET  /health           - Health check")
    print("=" * 60)
    print("\n🌐 Starting server on http://localhost:8000")
    print("📚 API Docs: http://localhost:8000/docs")
    print("=" * 60 + "\n")

    uvicorn.run(app, host="0.0.0.0", port=8000, reload=True)
