from pptx import Presentation
import os
import json # JSON डेटा को हैंडल करने के लिए यह लाइब्रेरी जोड़ी गई है

# MsoPlaceholderType values के लिए एक स्थानीय लुकअप डिक्शनरी (A local lookup dictionary for MsoPlaceholderType values)
# यह 'ImportError' से बचने के लिए सीधे pptx.enum.shapes को इम्पोर्ट करने से बचता है।
PLACEHOLDER_TYPES = {
    1: "TITLE (शीर्षक)",
    2: "BODY (मुख्य टेक्स्ट)",
    3: "CENTER_TITLE (केंद्र शीर्षक)",
    4: "SUBTITLE (उपशीर्षक)",
    5: "DATETIME (दिनांक)",
    6: "SLIDE_NUMBER (स्लाइड संख्या)",
    7: "FOOTER (फुटर)",
    8: "CONTENT (मीडिया/कंटेंट)",
    9: "PICTURE (चित्र)",
    10: "CHART (चार्ट)",
    11: "TABLE (तालिका)",
    12: "CLIP_ART (क्लिप आर्ट)",
    13: "DIAGRAM (डायग्राम)",
    14: "MEDIA_CLIP (मीडिया क्लिप)",
    15: "ORG_CHART (संगठन चार्ट)",
}

# इनपुट फ़ाइल का पथ (Path of the input file)
input_path = "templates/template-1.pptx"
# आउटपुट JSON फ़ाइल का पथ
output_path = "output_data.json"

def extract_placeholder_text(pptx_path):
    """
    दिए गए PPTX फ़ाइल से सभी प्लेसहोल्डर डेटा को निकालता है और JSON फ़ाइल में सेव करता है।
    """
    print(f"फाइल को पढ़ने की कोशिश कर रहा हूँ: {pptx_path}")
    
    # सुनिश्चित करें कि फ़ाइल मौजूद है (Ensure the file exists)
    if not os.path.exists(pptx_path):
        print(f"त्रुटि: फ़ाइल '{pptx_path}' नहीं मिली।")
        print("कृपया सुनिश्चित करें कि 'templates' फ़ोल्डर में 'template-1.pptx' फ़ाइल मौजूद है।")
        return

    try:
        # प्रेजेंटेशन ऑब्जेक्ट को लोड करें (Load the Presentation object)
        prs = Presentation(pptx_path)
        print("फाइल सफलतापूर्वक लोड हो गई। प्लेसहोल्डर की खोज कर रहा हूँ...")

        # सभी डेटा को स्टोर करने के लिए एक लिस्ट (List to store all data)
        extracted_data = []

        # हर स्लाइड पर Iterate करें (Iterate over every slide)
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": slide_idx + 1,
                "placeholders": []
            }
            
            # स्लाइड के सभी शेप्स (shapes) पर Iterate करें (Iterate over all shapes in the slide)
            for shape_idx, shape in enumerate(slide.shapes):
                
                # जांचें कि क्या शेप एक प्लेसहोल्डर है (Check if the shape is a placeholder)
                if shape.is_placeholder:
                    
                    placeholder_type_value = shape.placeholder_format.type
                    
                    # PLACEHOLDER_TYPES डिक्शनरी से नाम प्राप्त करें (Get the name from the PLACEHOLDER_TYPES dictionary)
                    placeholder_type_name = PLACEHOLDER_TYPES.get(
                        placeholder_type_value, 
                        f"UNKNOWN ({placeholder_type_value})"
                    )

                    placeholder_info = {
                        "index_in_slide": shape_idx + 1,
                        "type_id": placeholder_type_value,
                        "type_name": placeholder_type_name
                    }

                    # जांचें कि क्या प्लेसहोल्डर में टेक्स्ट फ्रेम है (Check if the placeholder has a text frame)
                    if shape.has_text_frame:
                        # टेक्स्ट निकालें (Extract the text)
                        text = shape.text.strip()
                        placeholder_info["text_content"] = text
                    else:
                        placeholder_info["text_content"] = "(Non-Text Placeholder)"
                    
                    slide_data["placeholders"].append(placeholder_info)
            
            # यदि स्लाइड में कोई प्लेसहोल्डर मिला, तो उसे मुख्य डेटा लिस्ट में जोड़ें
            if slide_data["placeholders"]:
                extracted_data.append(slide_data)

        # ------------------------------------------------------------------
        # JSON में डेटा सेव करें (Save data to JSON)
        print("\n--- डेटा प्रोसेसिंग पूरी हुई ---")
        with open(output_path, 'w', encoding='utf-8') as f:
            # json.dump का उपयोग करके डेटा को इंडेंटेशन के साथ JSON फॉर्मेट में सेव करें
            json.dump(extracted_data, f, ensure_ascii=False, indent=4)
        
        print(f"प्लेसहोल्डर डेटा सफलतापूर्वक JSON फॉर्मेट में सेव कर दिया गया है।")
        print(f"आउटपुट फ़ाइल पथ: {output_path}")
        # ------------------------------------------------------------------

    except Exception as e:
        print(f"\nफ़ाइल पढ़ते समय एक अंतिम त्रुटि हुई: {e}")

if __name__ == "__main__":
    # स्क्रिप्ट चलाएं (Run the script)
    extract_placeholder_text(input_path)